"""FastAPI application -- REST API + WebSocket for the OBD SuperStar Agent."""

from __future__ import annotations

import asyncio
import json
import logging
import uuid
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Optional

from fastapi import FastAPI, File, Form, Request, UploadFile, WebSocket, WebSocketDisconnect
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from starlette.middleware.base import BaseHTTPMiddleware

from backend.auth import (
    auth_enabled,
    auth_middleware,
    create_token,
    get_token_from_websocket,
    verify_token,
    authenticate_user,
)
from backend.config import OUTPUTS_DIR
from backend.database import (
    init_db,
    save_campaign,
    list_campaigns,
    get_campaign,
    delete_campaign,
    save_comment,
    list_comments,
    delete_comment,
)
from backend.collaboration import (
    register_user,
    unregister_user,
    update_user_activity,
    get_online_users,
    get_users_viewing_campaign,
    get_or_create_room,
    cleanup_room,
    record_activity,
    get_recent_activity,
    broadcast_to_all,
)
from backend.orchestrator import PipelineOrchestrator

# ── Logging ──
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(name)s] %(levelname)s: %(message)s",
)
logger = logging.getLogger(__name__)

# ── App ──
app = FastAPI(
    title="OBD SuperStar Agent",
    description="Multi-agent AI system for generating OBD promotional scripts and audio",
    version="1.0.0",
)

# CORS -- allow local dev and production frontend
import os
_allowed_origins = os.getenv("ALLOWED_ORIGINS", "http://localhost:3000,http://127.0.0.1:3000,http://localhost:8000").split(",")
# Strip whitespace from origins
_allowed_origins = [origin.strip() for origin in _allowed_origins if origin.strip()]

app.add_middleware(
    CORSMiddleware,
    allow_origins=_allowed_origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Auth middleware -- active when Supabase is configured or LOGIN_USERNAME/PASSWORD env vars are set
app.add_middleware(BaseHTTPMiddleware, dispatch=auth_middleware)

# Serve generated audio files
OUTPUTS_DIR.mkdir(exist_ok=True)
app.mount("/outputs", StaticFiles(directory=str(OUTPUTS_DIR)), name="outputs")

# Initialize campaign database
init_db()

# ── In-memory session store ──
sessions: dict[str, dict[str, Any]] = {}


# ── Background pipeline tracker ──

@dataclass
class PipelineState:
    """Tracks a running or completed pipeline."""
    session_id: str
    status: str = "running"  # running | done | error
    progress_log: list[dict[str, Any]] = field(default_factory=list)
    result: Optional[dict[str, Any]] = None
    error_message: str = ""
    task: Optional[asyncio.Task] = None  # type: ignore[type-arg]
    subscribers: list[WebSocket] = field(default_factory=list)

pipelines: dict[str, PipelineState] = {}


async def _run_pipeline_bg(
    state: PipelineState,
    product_text: str,
    country: str,
    telco: str,
    language: Optional[str],
    provider: Optional[str],
    tts_engine: Optional[str] = None,
) -> None:
    """Run the pipeline as a background task, storing progress in state."""

    async def on_progress(agent: str, status: str, data: dict[str, Any]) -> None:
        msg = {
            "agent": agent,
            "status": status,
            "message": data.get("message", ""),
            "data": {
                k: v for k, v in data.items()
                if k != "message" and _is_json_serializable(v)
            },
        }
        state.progress_log.append(msg)
        # Broadcast to any connected WebSocket subscribers
        dead: list[WebSocket] = []
        for ws in state.subscribers:
            try:
                await ws.send_json(msg)
            except Exception:
                dead.append(ws)
        for ws in dead:
            state.subscribers.remove(ws)

    try:
        orchestrator = PipelineOrchestrator(
            provider=provider,
            on_progress=on_progress,
        )
        result = await orchestrator.run(
            product_text=product_text,
            country=country,
            telco=telco,
            language=language,
            tts_engine=tts_engine,
        )
        session_id = result.get("session_id", state.session_id)
        result["country"] = country
        result["telco"] = telco
        result["language"] = language or ""
        result["tts_engine_choice"] = tts_engine or ""
        sessions[session_id] = result
        # Also store under the pipeline's session_id so the frontend can find it
        if state.session_id != session_id:
            sessions[state.session_id] = result
        state.result = _make_serializable(result)
        state.status = "error" if "error" in result else "done"

        # Notify subscribers of completion
        done_msg = {
            "agent": "Pipeline",
            "status": state.status,
            "message": result.get("error", "Pipeline complete"),
            "session_id": session_id,
            "result": state.result,
        }
        state.progress_log.append(done_msg)
        dead = []
        for ws in state.subscribers:
            try:
                await ws.send_json(done_msg)
            except Exception:
                dead.append(ws)
        for ws in dead:
            state.subscribers.remove(ws)

    except Exception as e:
        logger.exception(f"Background pipeline error: {e}")
        state.status = "error"
        state.error_message = str(e)
        err_msg = {
            "agent": "Pipeline",
            "status": "error",
            "message": str(e),
        }
        state.progress_log.append(err_msg)
        for ws in list(state.subscribers):
            try:
                await ws.send_json(err_msg)
            except Exception:
                pass


# ── REST Endpoints ──


@app.get("/api/health")
async def health_check():
    """Health check endpoint."""
    return {"status": "ok", "service": "OBD SuperStar Agent"}


@app.post("/api/auth/login")
async def login(request: Request):
    """Authenticate with username/password and receive a JWT cookie."""
    if not auth_enabled():
        return {"message": "Auth not enabled", "authenticated": True}

    body = await request.json()
    username = body.get("username", "")
    password = body.get("password", "")

    user = authenticate_user(username, password)
    if user:
        token = create_token(user)
        response = JSONResponse(content={
            "message": "Login successful",
            "authenticated": True,
            "username": username,
            "role": user.get("role"),
            "team": user.get("team"),
        })
        response.set_cookie(
            key="obd_token",
            value=token,
            httponly=True,
            samesite="lax",
            secure=request.url.scheme == "https",
            max_age=72 * 3600,  # 3 days
        )
        return response
    else:
        return JSONResponse(
            status_code=401,
            content={"error": "Invalid username or password"},
        )


@app.get("/api/auth/me")
async def auth_me(request: Request):
    """Check if the current user is authenticated."""
    if not auth_enabled():
        return {"authenticated": True, "auth_enabled": False, "username": "local", "role": "admin"}

    from backend.auth import get_token_from_request
    token = get_token_from_request(request)
    if token:
        payload = verify_token(token)
        if payload:
            return {
                "authenticated": True, 
                "auth_enabled": True, 
                "username": payload.get("sub"),
                "role": payload.get("role"),
                "team": payload.get("team"),
            }

    return JSONResponse(
        status_code=401,
        content={"authenticated": False, "auth_enabled": True},
    )


@app.post("/api/auth/logout")
async def logout():
    """Clear the auth cookie."""
    response = JSONResponse(content={"message": "Logged out"})
    response.delete_cookie("obd_token")
    return response


# ── Admin Endpoints ──

@app.get("/api/admin/users")
async def list_users():
    """List all users (Admin only)."""
    from backend.database import supabase
    if not supabase:
        return JSONResponse(status_code=500, content={"error": "Supabase not configured"})
        
    try:
        # Don't return password hashes to the frontend
        res = supabase.table("users").select("id, username, email, role, team, is_active, created_at").order("created_at").execute()
        return {"users": res.data}
    except Exception as e:
        logger.error(f"Failed to list users: {e}")
        return JSONResponse(status_code=500, content={"error": "Database error"})

@app.post("/api/admin/users")
async def create_user(request: Request):
    """Create a new user (Admin only)."""
    from backend.database import supabase
    import bcrypt
    
    if not supabase:
        return {"error": "Supabase not configured"}
        
    body = await request.json()
    username = body.get("username", "").strip()
    email = body.get("email", "").strip()
    password = body.get("password", "")
    role = body.get("role", "member")
    team = body.get("team", "default").strip()
    
    if not username or not email or not password:
        return JSONResponse(status_code=400, content={"error": "Missing required fields"})
        
    salt = bcrypt.gensalt()
    password_hash = bcrypt.hashpw(password.encode('utf-8'), salt).decode('utf-8')
    
    try:
        user_data = {
            "username": username,
            "email": email,
            "password_hash": password_hash,
            "role": role,
            "team": team,
            "is_active": True
        }
        res = supabase.table("users").insert(user_data).execute()
        
        # Log action
        admin_username = getattr(request.state, "username", "unknown")
        supabase.table("audit_log").insert({
            "admin_username": admin_username,
            "action": "create_user",
            "target_user": username,
            "details": {"role": role, "team": team}
        }).execute()
        
        # Don't return password hash
        new_user = res.data[0].copy()
        new_user.pop("password_hash", None)
        return {"user": new_user}
    except Exception as e:
        logger.error(f"Failed to create user: {e}")
        return JSONResponse(status_code=500, content={"error": "Failed to create user. Username or email might already exist."})

@app.put("/api/admin/users/{user_id}")
async def update_user(user_id: str, request: Request):
    """Update a user (Admin only)."""
    from backend.database import supabase
    import bcrypt
    
    if not supabase:
        return {"error": "Supabase not configured"}
        
    body = await request.json()
    updates = {}
    
    if "role" in body: updates["role"] = body["role"]
    if "team" in body: updates["team"] = body["team"].strip()
    if "is_active" in body: updates["is_active"] = body["is_active"]
    if "password" in body and body["password"]:
        salt = bcrypt.gensalt()
        updates["password_hash"] = bcrypt.hashpw(body["password"].encode('utf-8'), salt).decode('utf-8')
        
    if not updates:
        return {"message": "No updates provided"}
        
    try:
        res = supabase.table("users").update(updates).eq("id", user_id).execute()
        
        # Log action
        admin_username = getattr(request.state, "username", "unknown")
        target_username = res.data[0].get("username") if res.data else user_id
        
        log_details = updates.copy()
        log_details.pop("password_hash", None)
        
        supabase.table("audit_log").insert({
            "admin_username": admin_username,
            "action": "update_user",
            "target_user": target_username,
            "details": log_details
        }).execute()
        
        updated_user = res.data[0].copy() if res.data else {}
        updated_user.pop("password_hash", None)
        return {"user": updated_user}
    except Exception as e:
        logger.error(f"Failed to update user: {e}")
        return JSONResponse(status_code=500, content={"error": "Database error"})

@app.delete("/api/admin/users/{user_id}")
async def deactivate_user(user_id: str, request: Request):
    """Deactivate a user (Admin only, soft delete)."""
    from backend.database import supabase
    if not supabase:
        return {"error": "Supabase not configured"}
        
    try:
        res = supabase.table("users").update({"is_active": False}).eq("id", user_id).execute()
        
        # Log action
        admin_username = getattr(request.state, "username", "unknown")
        target_username = res.data[0].get("username") if res.data else user_id
        
        supabase.table("audit_log").insert({
            "admin_username": admin_username,
            "action": "deactivate_user",
            "target_user": target_username,
            "details": {}
        }).execute()
        
        return {"message": "User deactivated"}
    except Exception as e:
        logger.error(f"Failed to deactivate user: {e}")
        return JSONResponse(status_code=500, content={"error": "Database error"})


@app.post("/api/admin/cleanup")
async def admin_cleanup(request: Request):
    """Delete local output directories older than N hours. Admin only."""
    body = await request.json() if request.headers.get("content-type", "").startswith("application/json") else {}
    max_age_hours = body.get("max_age_hours", 2)
    deleted, freed = _cleanup_outputs(max_age_hours)
    return {
        "deleted_sessions": deleted,
        "freed_mb": round(freed / (1024 * 1024), 1),
    }


def _cleanup_outputs(max_age_hours: float = 2) -> tuple[int, int]:
    """Remove session output dirs older than *max_age_hours*. Returns (count, bytes)."""
    import time
    cutoff = time.time() - (max_age_hours * 3600)
    deleted = 0
    freed = 0
    if not OUTPUTS_DIR.exists():
        return 0, 0
    for child in OUTPUTS_DIR.iterdir():
        if not child.is_dir() or child.name.startswith("_"):
            continue
        try:
            mtime = max(f.stat().st_mtime for f in child.rglob("*") if f.is_file()) if any(child.rglob("*")) else child.stat().st_mtime
        except (StopIteration, OSError):
            mtime = child.stat().st_mtime
        if mtime < cutoff:
            size = sum(f.stat().st_size for f in child.rglob("*") if f.is_file())
            import shutil
            shutil.rmtree(child, ignore_errors=True)
            deleted += 1
            freed += size
    if deleted:
        logger.info(f"Cleanup: removed {deleted} session dirs, freed {freed / (1024*1024):.1f} MB")
    return deleted, freed


async def _periodic_cleanup():
    """Background task: clean up old outputs every hour."""
    while True:
        await asyncio.sleep(3600)
        try:
            _cleanup_outputs(max_age_hours=2)
        except Exception as e:
            logger.error(f"Periodic cleanup failed: {e}")


@app.on_event("startup")
async def _on_startup():
    _cleanup_outputs(max_age_hours=24)
    asyncio.create_task(_periodic_cleanup())


# ── File Upload / Text Extraction ──


@app.post("/api/upload/extract-text")
async def extract_text_from_file(file: UploadFile = File(...)):
    """Extract text content from uploaded documents (PDF, DOCX, PPTX, TXT, MD)."""
    if not file.filename:
        return JSONResponse(status_code=400, content={"error": "No file provided"})

    ext = Path(file.filename).suffix.lower()
    content = await file.read()

    try:
        if ext == ".pdf":
            import pdfplumber
            import io
            text_parts = []
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
            text = "\n\n".join(text_parts)

        elif ext in (".doc", ".docx"):
            import docx
            import io
            doc = docx.Document(io.BytesIO(content))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        elif ext == ".pptx":
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(content))
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            text = "\n\n".join(text_parts)

        elif ext == ".csv":
            text = content.decode("utf-8", errors="replace")

        elif ext in (".xlsx", ".xls"):
            import io
            try:
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
                text_parts = []
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        row_text = ", ".join(str(cell) for cell in row if cell is not None)
                        if row_text.strip():
                            text_parts.append(row_text)
                text = "\n".join(text_parts)
                wb.close()
            except ImportError:
                text = content.decode("utf-8", errors="replace")

        elif ext in (".txt", ".md", ".json", ".rtf"):
            text = content.decode("utf-8", errors="replace")

        else:
            return JSONResponse(
                status_code=400,
                content={"error": f"Unsupported file type: {ext}. Supported: PDF, DOCX, PPTX, TXT, CSV, XLSX, MD, JSON"},
            )

        if not text.strip():
            return JSONResponse(
                status_code=400,
                content={"error": "Could not extract any text from the file"},
            )

        logger.info(f"Extracted {len(text)} chars from {file.filename} ({ext})")
        return {"text": text, "filename": file.filename, "chars": len(text)}

    except Exception as e:
        logger.error(f"File extraction failed for {file.filename}: {e}")
        return JSONResponse(
            status_code=500,
            content={"error": f"Failed to extract text: {str(e)}"},
        )


# ── Background Pipeline Endpoints ──


@app.post("/api/generate/start")
async def start_pipeline(request: Request):
    """Start the pipeline as a background task. Returns session_id immediately."""
    body = await request.json()
    product_text = body.get("product_text", "")
    country = body.get("country", "")
    telco = body.get("telco", "")
    language = body.get("language")
    provider = body.get("provider")
    tts_engine = body.get("tts_engine")

    if not product_text or not country or not telco:
        return JSONResponse(
            status_code=400,
            content={"error": "product_text, country, and telco are required"},
        )

    session_id = uuid.uuid4().hex[:8]
    state = PipelineState(session_id=session_id)
    pipelines[session_id] = state

    task = asyncio.create_task(
        _run_pipeline_bg(state, product_text, country, telco, language, provider, tts_engine)
    )
    state.task = task

    logger.info(f"Pipeline started in background: session_id={session_id}")
    return {"session_id": session_id, "status": "running"}


@app.get("/api/generate/{session_id}/status")
async def pipeline_status(session_id: str):
    """Get current pipeline status, progress log, and result if done."""
    state = pipelines.get(session_id)
    if not state:
        # Check if it's a completed session from the old flow
        if session_id in sessions:
            return {
                "session_id": session_id,
                "status": "done",
                "progress": [],
                "result": _make_serializable(sessions[session_id]),
            }
        return JSONResponse(status_code=404, content={"error": "Pipeline not found"})

    resp: dict[str, Any] = {
        "session_id": session_id,
        "status": state.status,
        "progress": state.progress_log,
    }
    if state.result:
        resp["result"] = state.result
    if state.error_message:
        resp["error"] = state.error_message
    return resp


@app.post("/api/generate")
async def generate_scripts(
    product_file: Optional[UploadFile] = File(None),
    product_text: str = Form(""),
    country: str = Form(...),
    telco: str = Form(...),
    language: str = Form(""),
    provider: str = Form(""),
):
    """Start the OBD script generation pipeline (non-WebSocket version).

    Accepts product documentation as either a file upload or text input.
    Returns the full pipeline result synchronously.
    """
    # Get product text from file or form field
    if product_file:
        content = await product_file.read()
        doc_text = content.decode("utf-8", errors="replace")
    elif product_text:
        doc_text = product_text
    else:
        return JSONResponse(
            status_code=400,
            content={"error": "Either product_file or product_text is required"},
        )

    orchestrator = PipelineOrchestrator(provider=provider or None)

    result = await orchestrator.run(
        product_text=doc_text,
        country=country,
        telco=telco,
        language=language or None,
    )

    session_id = result.get("session_id", "unknown")
    result["country"] = country
    result["telco"] = telco
    result["language"] = language or ""
    sessions[session_id] = result

    return result


@app.get("/api/sessions/{session_id}")
async def get_session(session_id: str):
    """Retrieve results for a completed session."""
    if session_id in sessions:
        return sessions[session_id]
    return JSONResponse(status_code=404, content={"error": "Session not found"})


@app.get("/api/sessions/{session_id}/audio")
async def list_audio_files(session_id: str):
    """List all generated audio files for a session."""
    # First check if the session exists in memory to get public_urls
    if session_id in sessions and "audio" in sessions[session_id]:
        audio_files = sessions[session_id]["audio"].get("audio_files", [])
        if audio_files:
            files = []
            for file_info in audio_files:
                if "error" not in file_info:
                    files.append({
                        "name": file_info.get("file_name"),
                        "size_bytes": file_info.get("file_size_bytes"),
                        "url": file_info.get("public_url") or f"/outputs/{session_id}/{file_info.get('file_name')}",
                    })
            if files:
                return {"session_id": session_id, "files": files}
                
    # Fallback to local files if not in memory or no audio data
    session_dir = OUTPUTS_DIR / session_id
    if not session_dir.exists():
        return JSONResponse(status_code=404, content={"error": "Session not found"})

    files = []
    for f in sorted(list(session_dir.glob("*.mp3")) + list(session_dir.glob("*.wav"))):
        files.append({
            "name": f.name,
            "size_bytes": f.stat().st_size,
            "url": f"/outputs/{session_id}/{f.name}",
        })
    return {"session_id": session_id, "files": files}


@app.get("/api/audio/{session_id}/{filename}")
async def download_audio(session_id: str, filename: str, fmt: str = "mp3"):
    """Download a specific audio file.

    If the file was uploaded to Supabase Storage, redirect to the CDN URL.
    Otherwise serve from local disk (local dev fallback).
    """
    from starlette.responses import RedirectResponse

    public_url = _find_public_url(session_id, filename)
    if public_url:
        return RedirectResponse(url=public_url, status_code=302)

    file_path = OUTPUTS_DIR / session_id / filename
    if not file_path.exists():
        wav_alt = file_path.parent / (file_path.stem + ".wav")
        mp3_alt = file_path.parent / (file_path.stem + ".mp3")
        if wav_alt.exists():
            file_path = wav_alt
        elif mp3_alt.exists():
            file_path = mp3_alt
        else:
            return JSONResponse(status_code=404, content={"error": "File not found"})

    actual_ext = file_path.suffix.lower()

    if fmt == "wav" and actual_ext == ".mp3":
        try:
            from pydub import AudioSegment
            import io as _io
            seg = AudioSegment.from_mp3(str(file_path))
            buf = _io.BytesIO()
            seg.export(buf, format="wav")
            buf.seek(0)
            wav_name = file_path.stem + ".wav"
            from fastapi.responses import StreamingResponse
            return StreamingResponse(
                buf,
                media_type="audio/wav",
                headers={"Content-Disposition": f"attachment; filename={wav_name}"},
            )
        except Exception as e:
            logger.error("WAV conversion failed: %s", e)
            return JSONResponse(status_code=500, content={"error": "WAV conversion failed"})

    media_type = "audio/wav" if actual_ext == ".wav" else "audio/mpeg"
    return FileResponse(
        path=str(file_path),
        media_type=media_type,
        filename=file_path.name,
    )


def _find_public_url(session_id: str, filename: str) -> str | None:
    """Look up a Supabase public_url for an audio file from in-memory session data."""
    result = sessions.get(session_id)
    if not result:
        for pstate in pipelines.values():
            r = pstate.result
            if r and (r.get("session_id") == session_id):
                result = r
                break
    if not result:
        return None
    audio = result.get("audio", {})
    stem = Path(filename).stem
    for af in audio.get("audio_files", []):
        af_stem = Path(af.get("file_name", "")).stem
        if af_stem == stem and af.get("public_url"):
            return af["public_url"]
    return None


@app.get("/api/sessions/{session_id}/scripts")
async def download_scripts(session_id: str, fmt: str = "json", variant_id: Optional[int] = None):
    """Download scripts for a session as JSON or plain text. Optionally filter by variant_id."""
    if session_id not in sessions:
        # Fallback to DB if not in memory
        from backend.database import get_campaign
        campaign = get_campaign(session_id)
        if not campaign or not campaign.get("result"):
             return JSONResponse(status_code=404, content={"error": "Session not found"})
        result = campaign["result"]
    else:
        result = sessions[session_id]

    final_scripts = result.get("final_scripts", result.get("revised_scripts_round_1", result.get("initial_scripts", {})))
    scripts = final_scripts.get("scripts", [])

    if variant_id is not None:
        scripts = [s for s in scripts if s.get("variant_id") == variant_id]
        if not scripts:
            # Fallback to index if variant_id wasn't explicitly set
            if 0 <= variant_id - 1 < len(final_scripts.get("scripts", [])):
                scripts = [final_scripts["scripts"][variant_id - 1]]

    if not scripts:
        return JSONResponse(status_code=404, content={"error": "No scripts found in session"})

    filename_suffix = f"_v{variant_id}" if variant_id is not None else ""

    if fmt == "text":
        # Plain text format for easy reading / copy-paste
        lines = []
        for s in scripts:
            lines.append(f"{'='*60}")
            lines.append(f"VARIANT {s.get('variant_id', '?')}: {s.get('theme', '')}")
            lines.append(f"Language: {s.get('language', 'N/A')}  |  Words: {s.get('word_count', '?')}  |  ~{s.get('estimated_duration_seconds', '?')}s")
            lines.append(f"{'='*60}")
            lines.append(f"\n--- HOOK (0-5s) ---\n{s.get('hook', '')}")
            lines.append(f"\n--- BODY (5-23s) ---\n{s.get('body', '')}")
            lines.append(f"\n--- CTA (23-30s) ---\n{s.get('cta', '')}")
            lines.append(f"\n--- FULL SCRIPT ---\n{s.get('full_script', '')}")
            lines.append(f"\n--- FALLBACK 1 (Urgency) ---\n{s.get('fallback_1', '')}")
            lines.append(f"\n--- FALLBACK 2 (Psychology) ---\n{s.get('fallback_2', '')}")
            lines.append(f"\n--- POLITE CLOSURE ---\n{s.get('polite_closure', '')}")
            lines.append("")
        text_content = "\n".join(lines)
    from fastapi import Response
    if fmt == "text":
        return Response(
            content=text_content,
            media_type="text/plain",
            headers={"Content-Disposition": f"attachment; filename=scripts_{session_id}{filename_suffix}.txt"},
        )

    # Default: JSON
    return Response(
        content=json.dumps(final_scripts if variant_id is None else {"scripts": scripts}, indent=2),
        media_type="application/json",
        headers={"Content-Disposition": f"attachment; filename=scripts_{session_id}{filename_suffix}.json"},
    )


# ── Script Editing Endpoints ──


@app.put("/api/sessions/{session_id}/scripts/{variant_id}")
async def update_script(session_id: str, variant_id: int, request: Request):
    """Update a single script variant in a session (in-memory)."""
    if session_id not in sessions:
        return JSONResponse(status_code=404, content={"error": "Session not found"})

    body = await request.json()
    result = sessions[session_id]

    final_scripts = result.get("final_scripts") or result.get("revised_scripts_round_1") or result.get("initial_scripts")
    if not final_scripts:
        return JSONResponse(status_code=404, content={"error": "No scripts in session"})

    scripts = final_scripts.get("scripts", [])
    target = None
    for s in scripts:
        if s.get("variant_id") == variant_id:
            target = s
            break

    if not target:
        return JSONResponse(status_code=404, content={"error": f"Variant {variant_id} not found"})

    editable_fields = ("hook", "body", "cta", "full_script", "fallback_1", "fallback_2", "polite_closure")
    for field in editable_fields:
        if field in body:
            target[field] = body[field]

    word_count = len(target.get("full_script", "").split())
    target["word_count"] = word_count
    target["estimated_duration_seconds"] = round(word_count / 2.5, 1)

    logger.info(f"Updated script variant {variant_id} in session {session_id}")
    return {"status": "ok", "script": target}


@app.post("/api/sessions/{session_id}/regenerate-audio/{variant_id}")
async def regenerate_audio(session_id: str, variant_id: int, request: Request):
    """Regenerate audio files for a single script variant."""
    if session_id not in sessions:
        return JSONResponse(status_code=404, content={"error": "Session not found"})

    body = await request.json() if request.headers.get("content-type") == "application/json" else {}
    tts_engine_choice = body.get("tts_engine")

    result = sessions[session_id]
    country = result.get("country", "")
    language_val = result.get("language") or None

    final_scripts = result.get("final_scripts") or result.get("revised_scripts_round_1") or result.get("initial_scripts")
    if not final_scripts:
        return JSONResponse(status_code=404, content={"error": "No scripts in session"})

    target_script = None
    for s in final_scripts.get("scripts", []):
        if s.get("variant_id") == variant_id:
            target_script = s
            break

    if not target_script:
        return JSONResponse(status_code=404, content={"error": f"Variant {variant_id} not found"})

    voice_selection = result.get("voice_selection", {
        "selected_voice": {"voice_id": "", "name": "default"},
        "voice_settings": {},
    })

    single_scripts = {"scripts": [target_script]}

    from backend.agents import AudioProducerAgent
    producer = AudioProducerAgent()
    try:
        audio_result = await producer.run(
            scripts=single_scripts,
            voice_selection=voice_selection,
            session_id=session_id,
            country=country,
            language=language_val,
            tts_engine_override=tts_engine_choice,
        )

        new_files = audio_result.get("audio_files", [])

        # Update the session's audio data
        if "audio" not in result:
            result["audio"] = audio_result
        else:
            existing = result["audio"]
            existing["audio_files"] = [
                af for af in existing.get("audio_files", [])
                if af.get("variant_id") != variant_id
            ] + new_files
            existing["summary"]["total_generated"] = len([
                f for f in existing["audio_files"] if "error" not in f
            ])

        logger.info(f"Regenerated {len(new_files)} audio files for variant {variant_id}")
        return {
            "status": "ok",
            "audio_files": new_files,
            "tts_engine": audio_result.get("tts_engine"),
        }

    except Exception as e:
        logger.error(f"Audio regeneration failed for variant {variant_id}: {e}")
        return JSONResponse(status_code=500, content={"error": str(e)})


_audio_jobs: dict[str, dict[str, Any]] = {}


@app.post("/api/sessions/{session_id}/generate-full-audio")
async def generate_full_audio(session_id: str, request: Request):
    """Phase 2: kick off full audio generation as a background task.

    Returns a job_id immediately. Frontend polls /api/audio-jobs/{job_id}
    to check progress and retrieve the result -- avoids proxy timeout.
    """
    body = await request.json()
    voice_choices_raw = body.get("voice_choices", {})
    voice_choices = {int(k): int(v) for k, v in voice_choices_raw.items()}
    bgm_style = body.get("bgm_style", "upbeat")
    audio_format = body.get("audio_format", "mp3")
    tts_engine_choice = body.get("tts_engine")

    result = sessions.get(session_id)
    if not result:
        for pid, pstate in pipelines.items():
            if (pstate.result or {}).get("session_id") == session_id or pid == session_id:
                if pstate.result:
                    result = pstate.result
                    sessions[session_id] = result
                    break

    final_scripts = None
    if result:
        final_scripts = (
            result.get("final_scripts")
            or result.get("revised_scripts_round_1")
            or result.get("initial_scripts")
        )
    if not final_scripts:
        final_scripts = body.get("scripts")
    if not final_scripts:
        return JSONResponse(status_code=400, content={"error": "No scripts available. Please generate a new campaign."})

    voice_selection = (result or {}).get("voice_selection") or body.get("voice_selection") or {
        "selected_voice": {"voice_id": "", "name": "default"},
        "voice_settings": {},
    }
    country_val = body.get("country") or (result or {}).get("country", "")
    language_val = body.get("language") or (result or {}).get("language") or None

    job_id = uuid.uuid4().hex[:12]
    _audio_jobs[job_id] = {"status": "running", "session_id": session_id}

    async def _run():
        from backend.agents import AudioProducerAgent
        producer = AudioProducerAgent()
        try:
            audio_result = await producer.run_final_audio(
                scripts=final_scripts,
                voice_selection=voice_selection,
                voice_choices=voice_choices,
                session_id=session_id,
                country=country_val,
                language=language_val,
                tts_engine_override=tts_engine_choice or (result or {}).get("tts_engine_choice") or None,
                bgm_style=bgm_style,
                audio_format=audio_format,
            )
            if result:
                result["audio"] = audio_result
            else:
                sessions[session_id] = {"audio": audio_result, "session_id": session_id}
            logger.info(
                f"Full audio generated for session {session_id}: "
                f"{audio_result.get('summary', {}).get('total_generated', 0)} files"
            )
            _audio_jobs[job_id] = {
                "status": "done",
                "session_id": session_id,
                "audio": _make_serializable(audio_result),
            }
        except Exception as e:
            logger.error(f"Full audio generation failed for session {session_id}: {e}")
            _audio_jobs[job_id] = {"status": "error", "session_id": session_id, "error": str(e)}

    asyncio.create_task(_run())
    return {"status": "accepted", "job_id": job_id}


@app.get("/api/audio-jobs/{job_id}")
async def get_audio_job(job_id: str):
    """Poll for the status of a background audio generation job."""
    job = _audio_jobs.get(job_id)
    if not job:
        return JSONResponse(status_code=404, content={"error": "Job not found"})
    return job


# ── BGM Preview ──

_bgm_cache: dict[str, Path] = {}


@app.get("/api/bgm-preview/{style}")
async def bgm_preview(style: str):
    """Return a short (~8s) BGM-only MP3 sample for the given style."""
    from backend.agents.audio_producer import BGM_GENERATORS, _mix_voice_with_music
    import io as _io

    if style not in BGM_GENERATORS:
        return JSONResponse(status_code=400, content={"error": f"Unknown style: {style}"})

    if style in _bgm_cache and _bgm_cache[style].exists():
        return FileResponse(str(_bgm_cache[style]), media_type="audio/mpeg")

    gen = BGM_GENERATORS[style]
    wav_bytes = gen(8000)

    try:
        from pydub import AudioSegment
        seg = AudioSegment.from_wav(_io.BytesIO(wav_bytes))
        seg = seg - 10  # louder for preview (standalone, no voice)
        change_db = -16.0 - seg.dBFS
        seg = seg.apply_gain(change_db)
        preview_dir = OUTPUTS_DIR / "_bgm_previews"
        preview_dir.mkdir(exist_ok=True)
        out_path = preview_dir / f"{style}.mp3"
        seg.export(str(out_path), format="mp3", bitrate="192k")
        _bgm_cache[style] = out_path
        return FileResponse(str(out_path), media_type="audio/mpeg")
    except Exception as e:
        logger.error(f"BGM preview generation failed: {e}")
        return JSONResponse(status_code=500, content={"error": str(e)})


# ── Campaign Endpoints ──


@app.post("/api/campaigns")
async def create_campaign(request: Request):
    """Save a pipeline result as a named campaign."""
    body = await request.json()
    session_id = body.get("session_id", "")
    name = body.get("name", "").strip()

    if not session_id or not name:
        return JSONResponse(
            status_code=400,
            content={"error": "session_id and name are required"},
        )

    if session_id not in sessions:
        return JSONResponse(
            status_code=404,
            content={"error": "Session not found. Generate a campaign first."},
        )

    result = sessions[session_id]

    username = getattr(request.state, "username", "local")
    user_payload = getattr(request.state, "user", {})
    team = user_payload.get("team", "default") if isinstance(user_payload, dict) else "default"

    country = result.get("country", "")
    telco = result.get("telco", "")
    language = result.get("language", "")

    campaign = save_campaign(
        campaign_id=session_id,
        name=name,
        created_by=username,
        country=country,
        telco=telco,
        language=language,
        result=result,
        team=team,
    )

    return campaign


@app.get("/api/campaigns")
async def get_campaigns():
    """List all saved campaigns."""
    campaigns = list_campaigns()
    return {"campaigns": campaigns}


@app.get("/api/campaigns/{campaign_id}")
async def get_campaign_detail(campaign_id: str):
    """Get full details of a saved campaign."""
    campaign = get_campaign(campaign_id)
    if not campaign:
        return JSONResponse(status_code=404, content={"error": "Campaign not found"})
    return campaign


@app.delete("/api/campaigns/{campaign_id}")
async def remove_campaign(campaign_id: str):
    """Delete a saved campaign."""
    deleted = delete_campaign(campaign_id)
    if not deleted:
        return JSONResponse(status_code=404, content={"error": "Campaign not found"})
    return {"message": "Campaign deleted"}


# ── Collaboration Endpoints ──


@app.get("/api/presence")
async def get_presence():
    """Get all currently online users."""
    return {"users": get_online_users()}


@app.get("/api/presence/{campaign_id}")
async def get_campaign_presence(campaign_id: str):
    """Get users currently viewing a specific campaign."""
    return {"users": get_users_viewing_campaign(campaign_id)}


@app.get("/api/activity")
async def get_activity(limit: int = 20):
    """Get recent activity feed."""
    return {"events": get_recent_activity(limit)}


@app.get("/api/campaigns/{campaign_id}/comments")
async def get_comments(campaign_id: str):
    """List all comments for a campaign."""
    comments = list_comments(campaign_id)
    return {"comments": comments}


@app.post("/api/campaigns/{campaign_id}/comments")
async def add_comment(campaign_id: str, request: Request):
    """Add a comment to a campaign."""
    body = await request.json()
    text = body.get("text", "").strip()
    username = getattr(request.state, "username", body.get("username", "local"))
    if not text:
        return JSONResponse(status_code=400, content={"error": "Comment text is required"})

    comment_id = str(uuid.uuid4())
    comment = save_comment(comment_id, campaign_id, username, text)

    # Record activity and broadcast to collaboration room
    campaign = get_campaign(campaign_id)
    campaign_name = campaign["name"] if campaign else "Unknown"
    event = record_activity(
        "comment_added", username, campaign_id, campaign_name, text[:100]
    )

    room = get_or_create_room(campaign_id)
    await room.broadcast({"type": "comment_added", "comment": comment})
    await broadcast_to_all({"type": "activity", "event": event})

    return comment


@app.delete("/api/campaigns/{campaign_id}/comments/{comment_id}")
async def remove_comment(campaign_id: str, comment_id: str):
    """Delete a comment."""
    deleted = delete_comment(comment_id)
    if not deleted:
        return JSONResponse(status_code=404, content={"error": "Comment not found"})

    room = get_or_create_room(campaign_id)
    await room.broadcast({"type": "comment_deleted", "comment_id": comment_id})

    return {"deleted": True}


@app.websocket("/ws/collaborate/{campaign_id}")
async def websocket_collaborate(ws: WebSocket, campaign_id: str):
    """WebSocket for real-time collaboration on a campaign.

    On connect: sends current presence for the campaign.
    Receives: heartbeats and user actions.
    Broadcasts: user joined/left, comments, presence updates.
    """
    await ws.accept()
    ws_id = str(uuid.uuid4())

    from backend.auth import get_token_from_websocket, verify_token, auth_enabled
    username = "local"
    if auth_enabled():
        token = get_token_from_websocket(ws)
        if token:
            payload = verify_token(token)
            if payload:
                username = payload.get("sub") or "user"

    user = register_user(ws_id, username, ws)
    update_user_activity(ws_id, campaign_id)
    room = get_or_create_room(campaign_id)
    room.add(ws_id, ws)

    await room.broadcast(
        {"type": "user_joined", "user": user.to_dict()}, exclude_ws_id=ws_id
    )

    # Send current state to the connecting user
    await ws.send_json({
        "type": "init",
        "users": get_users_viewing_campaign(campaign_id),
        "comments": list_comments(campaign_id),
    })

    try:
        while True:
            data = await ws.receive_json()
            msg_type = data.get("type", "")

            if msg_type == "heartbeat":
                update_user_activity(ws_id, campaign_id)

            elif msg_type == "typing":
                await room.broadcast(
                    {"type": "typing", "username": username},
                    exclude_ws_id=ws_id,
                )
    except Exception:
        pass
    finally:
        room.remove(ws_id)
        left_user = unregister_user(ws_id)
        cleanup_room(campaign_id)
        if left_user:
            await room.broadcast({"type": "user_left", "username": left_user.username})


# ── WebSocket Endpoints ──


@app.websocket("/ws/progress/{session_id}")
async def websocket_progress(ws: WebSocket, session_id: str):
    """Read-only WebSocket that streams progress for a background pipeline.

    On connect: sends all buffered progress messages (catch-up).
    Then streams new messages as they arrive.
    Disconnect does NOT stop the pipeline.
    """
    if auth_enabled():
        token = get_token_from_websocket(ws)
        if not token or not verify_token(token):
            await ws.close(code=4001, reason="Authentication required")
            return

    state = pipelines.get(session_id)
    if not state:
        await ws.close(code=4004, reason="Pipeline not found")
        return

    await ws.accept()
    logger.info(f"Progress WS connected for session {session_id}")

    # Send all buffered progress (catch-up)
    for msg in list(state.progress_log):
        try:
            await ws.send_json(msg)
        except Exception:
            return

    # If already finished, close after catch-up
    if state.status in ("done", "error"):
        try:
            await ws.close()
        except Exception:
            pass
        return

    # Subscribe for live updates
    state.subscribers.append(ws)
    try:
        # Keep connection alive until client disconnects or pipeline finishes
        while True:
            # We just wait for the client to disconnect; all sending is done
            # from _run_pipeline_bg via the subscribers list
            try:
                await ws.receive_text()
            except WebSocketDisconnect:
                break
    finally:
        if ws in state.subscribers:
            state.subscribers.remove(ws)
        logger.info(f"Progress WS disconnected for session {session_id}")


@app.websocket("/ws/generate")
async def websocket_generate(ws: WebSocket):
    """WebSocket endpoint for real-time pipeline execution with progress updates.

    Client sends a JSON message to start:
    {
        "product_text": "...",
        "country": "...",
        "telco": "...",
        "language": "..." (optional),
        "provider": "..." (optional)
    }

    Server streams progress updates:
    {
        "agent": "AgentName",
        "status": "started|completed|error",
        "message": "...",
        "data": {...} (optional)
    }

    Final message includes the complete result.
    """
    # WebSocket auth check
    if auth_enabled():
        token = get_token_from_websocket(ws)
        if not token or not verify_token(token):
            await ws.close(code=4001, reason="Authentication required")
            return

    await ws.accept()
    logger.info("WebSocket client connected")

    try:
        # Wait for the initial configuration message
        raw = await ws.receive_text()
        config = json.loads(raw)

        product_text = config.get("product_text", "")
        country = config.get("country", "")
        telco = config.get("telco", "")
        language = config.get("language")
        provider = config.get("provider")
        tts_engine = config.get("tts_engine")

        if not product_text or not country or not telco:
            await ws.send_json({
                "agent": "Pipeline",
                "status": "error",
                "message": "product_text, country, and telco are required",
            })
            await ws.close()
            return

        # Progress callback that sends updates via WebSocket
        async def on_progress(agent: str, status: str, data: dict[str, Any]) -> None:
            try:
                await ws.send_json({
                    "agent": agent,
                    "status": status,
                    "message": data.get("message", ""),
                    "data": {
                        k: v for k, v in data.items()
                        if k != "message" and _is_json_serializable(v)
                    },
                })
            except Exception as e:
                logger.warning(f"Failed to send progress update: {e}")

        orchestrator = PipelineOrchestrator(
            provider=provider,
            on_progress=on_progress,
        )

        result = await orchestrator.run(
            product_text=product_text,
            country=country,
            telco=telco,
            language=language,
            tts_engine=tts_engine,
        )

        session_id = result.get("session_id", "unknown")
        result["country"] = country
        result["telco"] = telco
        result["language"] = language or ""
        result["tts_engine_choice"] = tts_engine or ""
        sessions[session_id] = result

        # Send the final result -- distinguish success vs failure
        has_error = "error" in result
        await ws.send_json({
            "agent": "Pipeline",
            "status": "error" if has_error else "done",
            "message": result.get("error", "Pipeline complete"),
            "session_id": session_id,
            "result": _make_serializable(result),
        })

    except WebSocketDisconnect:
        logger.info("WebSocket client disconnected")
    except json.JSONDecodeError:
        await ws.send_json({
            "agent": "Pipeline",
            "status": "error",
            "message": "Invalid JSON received",
        })
    except Exception as e:
        logger.exception(f"WebSocket error: {e}")
        try:
            await ws.send_json({
                "agent": "Pipeline",
                "status": "error",
                "message": str(e),
            })
        except Exception:
            pass
    finally:
        try:
            await ws.close()
        except Exception:
            pass


def _is_json_serializable(value: Any) -> bool:
    """Check if a value is JSON serializable."""
    try:
        json.dumps(value)
        return True
    except (TypeError, ValueError):
        return False



def _make_serializable(obj: Any) -> Any:
    """Recursively make an object JSON serializable."""
    if isinstance(obj, dict):
        return {k: _make_serializable(v) for k, v in obj.items()}
    elif isinstance(obj, list):
        return [_make_serializable(v) for v in obj]
    elif isinstance(obj, Path):
        return str(obj)
    elif _is_json_serializable(obj):
        return obj
    else:
        return str(obj)
